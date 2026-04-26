"""Generate PowerPoint presentation for AIAA 2711 Overfitting & Underfitting."""
import io
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


from lxml import etree

FIGURES = os.path.join(os.path.dirname(__file__), "..", "report", "figures")
OUT = os.path.join(os.path.dirname(__file__), "presentation_v5.pptx")

# Colors
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
BG_SLIDE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2D, 0x6A, 0x9F)
ACCENT_LIGHT = RGBColor(0xE8, 0xF0, 0xF8)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_GRAY = RGBColor(0x55, 0x55, 0x55)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)


def render_latex(tex: str, fontsize: int = 18, dpi: int = 200) -> io.BytesIO:
    plt.rcParams.update({
        "text.usetex": True,
        "text.latex.preamble": r"\usepackage{amsmath,amssymb}",
    })
    fig, ax = plt.subplots(figsize=(0.01, 0.01))
    ax.axis("off")
    text = ax.text(0, 0, f"${tex}$", fontsize=fontsize,
                   color="black", ha="left", va="bottom",
                   transform=ax.transAxes)
    fig.patch.set_alpha(0)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                pad_inches=0.05, transparent=True)
    plt.close(fig)
    buf.seek(0)
    plt.rcParams.update({"text.usetex": False})
    return buf


def add_bg(slide, color=BG_SLIDE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT_DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_accent_bar(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def add_image(slide, path, left, top, width=None, height=None):
    if width and height:
        pic = slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        pic = slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        pic = slide.shapes.add_picture(path, left, top, height=height)
    else:
        pic = slide.shapes.add_picture(path, left, top)
    return pic


def add_latex_image(slide, tex, left, top, fontsize=18, dpi=200):
    buf = render_latex(tex, fontsize=fontsize, dpi=dpi)
    return slide.shapes.add_picture(buf, left, top)


def add_block(slide, left, top, width, height, title, body, title_color=WHITE,
              bg_color=ACCENT, body_color=TEXT_DARK):
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = title_color
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Body
    body_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.4),
                                        width, height - Inches(0.4))
    body_shape.fill.solid()
    body_shape.fill.fore_color.rgb = ACCENT_LIGHT
    body_shape.line.fill.background()
    tf2 = body_shape.text_frame
    tf2.word_wrap = True
    tf2.paragraphs[0].text = body
    tf2.paragraphs[0].font.size = Pt(13)
    tf2.paragraphs[0].font.color.rgb = body_color
    tf2.paragraphs[0].font.name = "Calibri"
    return shape, body_shape


def make_eigenvalue_diagram() -> io.BytesIO:
    """Visualize eigenvalue shift: adding alpha*I stabilizes X'X."""
    fig, ax = plt.subplots(figsize=(5.5, 2.4))

    lambdas = np.array([0.05, 0.4, 1.2, 2.3, 3.8])
    alpha = 1.0
    shifted = lambdas + alpha

    # Danger zone near zero
    ax.axvspan(-0.15, 0.15, alpha=0.15, color='#E74C3C')
    ax.text(0, -0.25, 'unstable\nzone', ha='center', fontsize=8,
            color='#E74C3C', fontstyle='italic')

    y_top, y_bot = 0.75, 0.25
    for i in range(len(lambdas)):
        col = '#E74C3C' if lambdas[i] < 0.1 else ('#F39C12' if lambdas[i] < 0.5 else '#2D6A9F')
        ax.plot(lambdas[i], y_top, 'o', markersize=11, color=col, zorder=5)
        ax.plot(shifted[i], y_bot, 'o', markersize=11, color='#2D6A9F', zorder=5)
        ax.annotate('', xy=(shifted[i], y_bot + 0.04),
                    xytext=(lambdas[i], y_top - 0.04),
                    arrowprops=dict(arrowstyle='->', color='#999999', lw=0.7))

    ax.text(-0.5, y_top, r"$X'X$", fontsize=11, fontweight='bold',
            ha='right', va='center')
    ax.text(-0.5, y_bot, r"$X'X+\alpha I$", fontsize=11, fontweight='bold',
            ha='right', va='center', color='#2D6A9F')

    ax.set_xlim(-0.8, 5.2)
    ax.set_ylim(-0.5, 1.1)
    ax.set_yticks([])
    ax.spines['top'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_xlabel('Eigenvalue', fontsize=10)
    ax.axhline(y=y_top, color='gray', lw=0.3)
    ax.axhline(y=y_bot, color='gray', lw=0.3)

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf


def make_ridge_constraint_contour() -> io.BytesIO:
    """Ridge: L2 constraint circle with loss function contours."""
    fig, ax = plt.subplots(figsize=(3.8, 3.8))

    w_ols = np.array([2.3, 1.5])
    t = 1.3
    Sigma = np.array([[0.7, 0.15], [0.15, 1.0]])

    w1 = np.linspace(-2, 3.5, 300)
    w2 = np.linspace(-2, 3.5, 300)
    W1, W2 = np.meshgrid(w1, w2)
    d1, d2 = W1 - w_ols[0], W2 - w_ols[1]
    Z = Sigma[0, 0] * d1**2 + 2 * Sigma[0, 1] * d1 * d2 + Sigma[1, 1] * d2**2

    ax.contour(W1, W2, Z, levels=8, colors='#2D6A9F', linewidths=0.8, alpha=0.5)

    # Constraint circle ||w||^2 <= t
    theta = np.linspace(0, 2 * np.pi, 200)
    ax.plot(t * np.cos(theta), t * np.sin(theta), color='#E86C00', linewidth=2.5)
    ax.fill(t * np.cos(theta), t * np.sin(theta), color='#E86C00', alpha=0.07)
    ax.text(t * 0.5, -t * 0.6, r'$||w||^2 \leq t$', fontsize=10,
            color='#E86C00', fontstyle='italic')

    # Ridge solution: tangent point on circle (numerical search)
    best_th, best_val = 0, float('inf')
    for th in np.linspace(0, 2 * np.pi, 2000):
        w = np.array([t * np.cos(th), t * np.sin(th)])
        d = w - w_ols
        val = d @ Sigma @ d
        if val < best_val:
            best_val = val
            best_th = th
    w_ridge = np.array([t * np.cos(best_th), t * np.sin(best_th)])

    ax.plot(*w_ols, 'x', color='#555', markersize=13, markeredgewidth=2.5,
            label=r'$w_{OLS}$')
    ax.plot(*w_ridge, 'o', color='#E86C00', markersize=10, markeredgewidth=2.5,
            markerfacecolor='white', zorder=5, label=r'$w_{Ridge}$')

    ax.set_xlabel(r'$w_1$', fontsize=11)
    ax.set_ylabel(r'$w_2$', fontsize=11)
    ax.set_aspect('equal')
    ax.legend(fontsize=9, loc='upper left')
    ax.set_xlim(-2, 3.2)
    ax.set_ylim(-2, 3)
    ax.axhline(y=0, color='gray', lw=0.3)
    ax.axvline(x=0, color='gray', lw=0.3)
    ax.grid(True, alpha=0.12)

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf


def make_l1_vs_l2_geometry() -> io.BytesIO:
    """Classic L2 (circle) vs L1 (diamond) constraint comparison."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 3.5))

    w_ols = np.array([2.5, 1.0])
    Sigma = np.array([[0.6, 0.2], [0.2, 1.2]])

    w1 = np.linspace(-1.5, 3.5, 300)
    w2 = np.linspace(-1.5, 2.5, 300)
    W1, W2 = np.meshgrid(w1, w2)
    d1, d2 = W1 - w_ols[0], W2 - w_ols[1]
    Z = Sigma[0, 0] * d1**2 + 2 * Sigma[0, 1] * d1 * d2 + Sigma[1, 1] * d2**2

    t = 1.2

    for ax, mode in [(ax1, 'L2'), (ax2, 'L1')]:
        ax.contour(W1, W2, Z, levels=10, colors='#2D6A9F', linewidths=0.7, alpha=0.4)

        if mode == 'L2':
            theta = np.linspace(0, 2 * np.pi, 200)
            ax.plot(t * np.cos(theta), t * np.sin(theta), color='#2D6A9F', linewidth=2.5)
            ax.fill(t * np.cos(theta), t * np.sin(theta), color='#2D6A9F', alpha=0.08)

            best_th, best_val = 0, float('inf')
            for th in np.linspace(0, 2 * np.pi, 2000):
                w = np.array([t * np.cos(th), t * np.sin(th)])
                d = w - w_ols
                val = d @ Sigma @ d
                if val < best_val:
                    best_val = val
                    best_th = th
            w_sol = np.array([t * np.cos(best_th), t * np.sin(best_th)])
            ax.plot(*w_sol, 'o', color='#2D6A9F', markersize=10, markeredgewidth=2.5,
                    markerfacecolor='white', zorder=5)
            ax.annotate(r'$w_1 \neq 0,\; w_2 \neq 0$', w_sol,
                        textcoords="offset points", xytext=(12, 12),
                        fontsize=9, color='#2D6A9F')
            ax.set_title('Ridge (L2): Circle', fontsize=13,
                         fontweight='bold', color='#2D6A9F')
        else:
            diamond_x = [t, 0, -t, 0, t]
            diamond_y = [0, t, 0, -t, 0]
            ax.plot(diamond_x, diamond_y, color='#E86C00', linewidth=2.5)
            ax.fill(diamond_x, diamond_y, color='#E86C00', alpha=0.08)

            # Lasso solution at corner — sparse!
            ax.plot(t, 0, 'o', color='#E86C00', markersize=10, markeredgewidth=2.5,
                    markerfacecolor='white', zorder=5)
            ax.annotate(r'$w_2 = 0$ !', (t, 0),
                        textcoords="offset points", xytext=(12, 12),
                        fontsize=10, color='#E86C00', fontweight='bold')
            ax.set_title('Lasso (L1): Diamond', fontsize=13,
                         fontweight='bold', color='#E86C00')

        ax.plot(*w_ols, 'x', color='#555', markersize=12, markeredgewidth=2.5)
        ax.set_xlabel(r'$w_1$', fontsize=11)
        ax.set_ylabel(r'$w_2$', fontsize=11)
        ax.set_aspect('equal')
        ax.set_xlim(-1.8, 3.5)
        ax.set_ylim(-1.8, 2.5)
        ax.axhline(y=0, color='gray', lw=0.3)
        ax.axvline(x=0, color='gray', lw=0.3)
        ax.grid(True, alpha=0.12)

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf


P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'


def _pn(tag):
    return f'{{{P_NS}}}{tag}'


def add_appear_animations(slide, click_groups):
    """Add click-to-appear animations.  click_groups: list of lists of shapes."""
    id_ctr = [1]
    def nid():
        v = id_ctr[0]; id_ctr[0] += 1; return str(v)

    for child in list(slide._element):
        if child.tag == _pn('timing'):
            slide._element.remove(child)

    timing = etree.SubElement(slide._element, _pn('timing'))
    tnLst = etree.SubElement(timing, _pn('tnLst'))
    root_par = etree.SubElement(tnLst, _pn('par'))
    root_cTn = etree.SubElement(root_par, _pn('cTn'))
    root_cTn.set('id', nid()); root_cTn.set('dur', 'indefinite')
    root_cTn.set('restart', 'never'); root_cTn.set('nodeType', 'tmRoot')
    root_child = etree.SubElement(root_cTn, _pn('childTnLst'))

    seq = etree.SubElement(root_child, _pn('seq'))
    seq.set('concurrent', '1'); seq.set('nextAc', 'seek')
    seq_cTn = etree.SubElement(seq, _pn('cTn'))
    seq_cTn.set('id', nid()); seq_cTn.set('dur', 'indefinite')
    seq_cTn.set('nodeType', 'mainSeq')
    seq_child = etree.SubElement(seq_cTn, _pn('childTnLst'))

    for group in click_groups:
        click_par = etree.SubElement(seq_child, _pn('par'))
        click_cTn = etree.SubElement(click_par, _pn('cTn'))
        click_cTn.set('id', nid()); click_cTn.set('fill', 'hold')
        sc = etree.SubElement(click_cTn, _pn('stCondLst'))
        etree.SubElement(sc, _pn('cond')).set('delay', '0')
        click_child = etree.SubElement(click_cTn, _pn('childTnLst'))

        for shape in group:
            sp = etree.SubElement(click_child, _pn('par'))
            sp_cTn = etree.SubElement(sp, _pn('cTn'))
            sp_cTn.set('id', nid()); sp_cTn.set('fill', 'hold')
            s = etree.SubElement(sp_cTn, _pn('stCondLst'))
            etree.SubElement(s, _pn('cond')).set('delay', '0')
            sp_child = etree.SubElement(sp_cTn, _pn('childTnLst'))

            p_set = etree.SubElement(sp_child, _pn('set'))
            cBhvr = etree.SubElement(p_set, _pn('cBhvr'))
            b_cTn = etree.SubElement(cBhvr, _pn('cTn'))
            b_cTn.set('id', nid()); b_cTn.set('dur', '1'); b_cTn.set('fill', 'hold')
            bs = etree.SubElement(b_cTn, _pn('stCondLst'))
            etree.SubElement(bs, _pn('cond')).set('delay', '0')

            tgt = etree.SubElement(cBhvr, _pn('tgtEl'))
            etree.SubElement(tgt, _pn('spTgt')).set('spid', str(shape.shape_id))

            anl = etree.SubElement(cBhvr, _pn('attrNameLst'))
            an = etree.SubElement(anl, _pn('attrName'))
            an.text = 'style.visibility'

            to = etree.SubElement(p_set, _pn('to'))
            etree.SubElement(to, _pn('strVal')).set('val', 'visible')

    prev = etree.SubElement(seq, _pn('prevCondLst'))
    pc = etree.SubElement(prev, _pn('cond'))
    pc.set('evt', 'onPrev'); pc.set('delay', '0')
    etree.SubElement(etree.SubElement(pc, _pn('tgtEl')), _pn('sldTgt'))

    nxt = etree.SubElement(seq, _pn('nextCondLst'))
    nc = etree.SubElement(nxt, _pn('cond'))
    nc.set('evt', 'onNext'); nc.set('delay', '0')
    etree.SubElement(etree.SubElement(nc, _pn('tgtEl')), _pn('sldTgt'))

    bldLst = etree.SubElement(timing, _pn('bldLst'))
    for group in click_groups:
        for shape in group:
            bp = etree.SubElement(bldLst, _pn('bldP'))
            bp.set('spid', str(shape.shape_id)); bp.set('grpId', '0')


# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]  # blank

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(2.5), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.5),
            "Overfitting and Underfitting\nin Machine Learning",
            font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
            "AIAA 2711 \u2014 Spring 2025",
            font_size=22, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
            "Member A  \u00b7  Member B  \u00b7  Member C",
            font_size=20, color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER)

# slide number placeholder
add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "1 / 14", font_size=10, color=RGBColor(0x66, 0x66, 0x66),
            alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "Welcome everyone. My name is [Speaker 1], and together with [Speaker 2] and "
    "[Speaker 3], we will present our study on overfitting and underfitting in "
    "machine learning. (10s)"
)


# ============================================================
# SLIDE 2: The Hook
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "The Hook: Which Model Do You Trust?",
            font_size=30, bold=True, color=ACCENT)

fig_path = os.path.join(FIGURES, "polynomial_fits.png")
add_image(slide, fig_path, Inches(0.8), Inches(1.2), width=Inches(11.5))

add_textbox(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.6),
            "\u2728 Vote by show of hands: degree 1, degree 4, or degree 15?",
            font_size=20, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

add_block(slide, Inches(3), Inches(6.3), Inches(7), Inches(0.8),
          "Interactive Question",
          "Which polynomial would you trust to predict a NEW data point?")

add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "2 / 14", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "[INTERACTION] Show 3 plots. Ask audience: 'Which polynomial would you trust to "
    "predict a new data point?' Quick hand raise for degree 1, 4, 15.\n\n"
    "'This is the fundamental tension in ML \u2014 too simple vs too complex. "
    "Today we derive the mathematics behind WHY this happens and HOW to fix it.' (50s)"
)


# ============================================================
# SLIDE 3: Formal Definitions
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Formal Definitions", font_size=30, bold=True, color=ACCENT)

# Training Error
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.4),
            "Training Error:", font_size=20, bold=True, color=TEXT_DARK)
add_latex_image(slide,
    r"E_{\mathrm{train}} = \frac{1}{N}\sum_{i=1}^{N} \ell(y_i, \hat{f}(x_i))",
    Inches(1.2), Inches(1.8), fontsize=20)

# Test Error
add_textbox(slide, Inches(0.8), Inches(3.0), Inches(5), Inches(0.4),
            "Test Error:", font_size=20, bold=True, color=TEXT_DARK)
add_latex_image(slide,
    r"E_{\mathrm{test}} = \mathbb{E}_{(x,y)\sim P}[\ell(y, \hat{f}(x))]",
    Inches(1.2), Inches(3.5), fontsize=20)

# Definitions box on right
add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
            "Underfitting", font_size=22, bold=True, color=ACCENT)
add_textbox(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(0.6),
            "E_train is HIGH \u2014 model too simple.\n"
            "Both training and test errors are large.",
            font_size=16, color=TEXT_GRAY)

add_textbox(slide, Inches(7), Inches(3.2), Inches(5.5), Inches(0.5),
            "Overfitting", font_size=22, bold=True, color=ORANGE)
add_textbox(slide, Inches(7), Inches(3.8), Inches(5.5), Inches(0.6),
            "E_train \u226a E_test \u2014 generalization gap.\n"
            "Model memorizes noise in training data.",
            font_size=16, color=TEXT_GRAY)

add_block(slide, Inches(2), Inches(5.3), Inches(9), Inches(0.85),
          "Key Message",
          "The generalization gap (E_test \u2212 E_train) is our diagnostic tool.")

add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "3 / 14", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "Define training and test error formally. Underfitting = high training error; "
    "overfitting = large gap between train and test. 'The generalization gap is our "
    "diagnostic.' (40s)"
)


# ============================================================
# SLIDE 4: The Key Question
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Why Does Test Error Follow a U-Shape?",
            font_size=30, bold=True, color=ACCENT)

# Left column
add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3),
    ["\u2022  Training error \u2193 monotonically with complexity",
     "\u2022  Test error is U-shaped",
     "\u2022  Minimum near degree 3\u20135"],
    font_size=18)

add_textbox(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(0.5),
            "The Answer:", font_size=22, bold=True, color=ORANGE)
add_textbox(slide, Inches(0.8), Inches(4.4), Inches(5.5), Inches(1),
            "Bias-Variance Decomposition\n(uses expectation & variance from Weeks 9\u201310)",
            font_size=18, color=TEXT_DARK)

# Right: figure
fig_path = os.path.join(FIGURES, "error_vs_complexity.png")
add_image(slide, fig_path, Inches(6.8), Inches(1.2), width=Inches(5.8))

add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "4 / 14", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "Training error always decreases, but test error has a U-shape. The sweet spot "
    "is around degree 3\u20135. We explain this via bias-variance decomposition. "
    "'This uses expectation and variance from Weeks 9 and 10.' (20s)"
)


# ============================================================
# SLIDE 5: Bias-Variance Setup
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Bias-Variance Decomposition: Setup",
            font_size=30, bold=True, color=ACCENT)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(4), Inches(0.4),
            "Data Model:", font_size=20, bold=True, color=TEXT_DARK)
add_latex_image(slide,
    r"y = f(x) + \epsilon, \quad \epsilon \sim \mathcal{N}(0, \sigma^2)",
    Inches(1.2), Inches(1.9), fontsize=22)

add_bullet_list(slide, Inches(0.8), Inches(3.0), Inches(5.5), Inches(1.5),
    ["f(x) : true underlying function",
     "\u0192\u0302(x) : model trained on dataset D",
     "\u03b5 : irreducible noise (Gaussian)"],
    font_size=17)

add_textbox(slide, Inches(0.8), Inches(4.8), Inches(5), Inches(0.4),
            "Our Goal \u2014 decompose:", font_size=20, bold=True, color=TEXT_DARK)
add_latex_image(slide,
    r"\mathbb{E}_D\left[(y - \hat{f}(x))^2\right]",
    Inches(1.5), Inches(5.4), fontsize=22)

add_block(slide, Inches(7), Inches(2.5), Inches(5.5), Inches(1.2),
          "Key Point",
          "The expectation is over ALL possible training datasets D. "
          "We ask: on average, across all training sets we could draw, "
          "how far off is our model?")

add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "5 / 14", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "Set up the bias-variance framework. y = f(x) + noise. We train f-hat on dataset D. "
    "The key insight: expectation is over ALL possible training sets D. (40s)"
)


# ============================================================
# SLIDE 6a: Bias-Variance Derivation — Steps (CORE)
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Bias-Variance Derivation: Step-by-Step \u2605",
            font_size=28, bold=True, color=ACCENT)

# --- Click 1: Step 1 full expansion ---
s1_title = add_textbox(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.4),
            "Step 1: Substitute y = f + \u03b5 and expand", font_size=17, bold=True, color=ORANGE)
s1_formula = add_latex_image(slide,
    r"\mathbb{E}_D[(y{-}\hat{f})^2] = \mathbb{E}_D[(f{-}\hat{f})^2] + \mathbb{E}_D[\epsilon^2] + 2\,\mathbb{E}_D[\epsilon(f{-}\hat{f})]",
    Inches(0.9), Inches(1.6), fontsize=15)

# --- Click 2: Cross-term analysis ---
s1_cross = add_textbox(slide, Inches(0.9), Inches(2.3), Inches(11), Inches(0.7),
            "Cross-term: \u03b5 is independent of \u0192\u0302, and f is fixed, so\n"
            "E[\u03b5(f \u2212 \u0192\u0302)] = E[\u03b5] \u00b7 E[f \u2212 \u0192\u0302] = 0 \u00b7 (\u2026) = 0    "
            "Also: E[\u03b5\u00b2] = \u03c3\u00b2",
            font_size=13, color=TEXT_GRAY)
s1_result_formula = add_latex_image(slide,
    r"\boxed{\mathbb{E}_D[(y{-}\hat{f})^2] = \mathbb{E}_D[(f{-}\hat{f})^2] + \sigma^2}",
    Inches(2.5), Inches(3.1), fontsize=16)

# --- Click 3: Step 2 expansion ---
s2_title = add_textbox(slide, Inches(0.8), Inches(3.9), Inches(6), Inches(0.4),
            "Step 2: Let \u0192\u0304 \u2261 E_D[\u0192\u0302], add-subtract \u0192\u0304",
            font_size=17, bold=True, color=ORANGE)
s2_formula = add_latex_image(slide,
    r"\mathbb{E}_D[(f{-}\hat{f})^2] = \underbrace{(f{-}\bar{f})^2}_{\text{constant}} + \mathbb{E}_D[(\hat{f}{-}\bar{f})^2] + 2(f{-}\bar{f})\underbrace{\mathbb{E}_D[\bar{f}{-}\hat{f}]}_{=\bar{f}-\bar{f}=0}",
    Inches(0.9), Inches(4.4), fontsize=14)

# --- Click 4: Cross-term explanation ---
s2_cross = add_textbox(slide, Inches(0.9), Inches(5.4), Inches(11), Inches(0.5),
            "Cross-term: E_D[\u0192\u0304 \u2212 \u0192\u0302] = \u0192\u0304 \u2212 E_D[\u0192\u0302] = \u0192\u0304 \u2212 \u0192\u0304 = 0  "
            "(by definition of \u0192\u0304)",
            font_size=13, color=TEXT_GRAY)

# Add click-to-appear animations (4 clicks)
add_appear_animations(slide, [
    [s1_title, s1_formula],                        # Click 1: Step 1 expansion
    [s1_cross, s1_result_formula],                 # Click 2: Cross-term + boxed result
    [s2_title, s2_formula],                        # Click 3: Step 2 expansion
    [s2_cross],                                     # Click 4: Step 2 cross-term
])

add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "6 / 14", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "[CORE SLIDE 6a - Step through EVERY term]\n"
    "Step 1: Substitute y=f+eps, expand square -> 3 terms. "
    "Cross-term: eps independent of f-hat, E[eps]=0, so cross-term=0. E[eps^2]=sigma^2.\n"
    "Step 2: Let f-bar = E_D[f-hat]. Add-subtract f-bar, expand -> 3 terms. "
    "Cross-term: E_D[f-bar - f-hat] = f-bar - f-bar = 0 by definition. (50s)"
)


# ============================================================
# SLIDE 6b: Bias-Variance Result + Monte Carlo Verification
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(7), Inches(0.7),
            "Bias-Variance: Result & Verification",
            font_size=28, bold=True, color=ACCENT)

# Final result box
result_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(1.1), Inches(7.5), Inches(1.6))
result_box.fill.solid()
result_box.fill.fore_color.rgb = RGBColor(0xFD, 0xF6, 0xE3)
result_box.line.color.rgb = ORANGE
result_box.line.width = Pt(2)

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(4), Inches(0.3),
            "Combining Steps 1 & 2:", font_size=16, bold=True, color=ORANGE)
add_latex_image(slide,
    r"\mathbb{E}_D[(y{-}\hat{f})^2] = \underbrace{(f - \bar{f})^2}_{\mathrm{Bias}^2} + \underbrace{\mathbb{E}_D[(\hat{f} - \bar{f})^2]}_{\mathrm{Variance}} + \underbrace{\sigma^2}_{\mathrm{Irreducible}}",
    Inches(0.8), Inches(1.7), fontsize=18)

# Interpretation
add_textbox(slide, Inches(0.5), Inches(3.0), Inches(7.5), Inches(0.4),
            "Interpretation:", font_size=18, bold=True, color=TEXT_DARK)

# Bias box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(3.5), Inches(3.7), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_LIGHT
shape.line.color.rgb = ACCENT
add_textbox(slide, Inches(0.7), Inches(3.6), Inches(3.3), Inches(0.3),
            "Bias\u00b2", font_size=15, bold=True, color=ACCENT)
add_textbox(slide, Inches(0.7), Inches(3.95), Inches(3.3), Inches(0.6),
            "Avg model's distance from truth\nHigh bias \u2192 underfitting",
            font_size=13, color=TEXT_DARK)

# Variance box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(4.4), Inches(3.5), Inches(3.7), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFD, 0xF6, 0xE3)
shape.line.color.rgb = ORANGE
add_textbox(slide, Inches(4.6), Inches(3.6), Inches(3.3), Inches(0.3),
            "Variance", font_size=15, bold=True, color=ORANGE)
add_textbox(slide, Inches(4.6), Inches(3.95), Inches(3.3), Inches(0.6),
            "Prediction fluctuation across datasets\nHigh variance \u2192 overfitting",
            font_size=13, color=TEXT_DARK)

# Tradeoff
add_textbox(slide, Inches(0.5), Inches(5.0), Inches(7.5), Inches(0.7),
            "Complexity \u2191 \u2192 Bias \u2193 Variance \u2191    |    "
            "Complexity \u2193 \u2192 Bias \u2191 Variance \u2193",
            font_size=14, bold=True, color=TEXT_DARK)

# Monte Carlo figure on right
fig_path = os.path.join(FIGURES, "bias_variance_tradeoff.png")
add_image(slide, fig_path, Inches(8.3), Inches(1.0), width=Inches(4.5))

add_textbox(slide, Inches(8.3), Inches(5.3), Inches(4.5), Inches(0.8),
            "Monte Carlo verification\n200 independent datasets",
            font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# sigma box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(5.8), Inches(7.5), Inches(0.7))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
add_textbox(slide, Inches(0.7), Inches(5.9), Inches(7), Inches(0.4),
            "\u03c3\u00b2 = irreducible noise floor \u2014 no model can beat this",
            font_size=13, color=TEXT_GRAY)

add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "7 / 14", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "[CORE SLIDE 6b] Present boxed result. Interpret: Bias = avg model distance from truth "
    "(underfitting). Variance = prediction fluctuation (overfitting). sigma^2 = noise floor.\n"
    "[POINT TO FIGURE] Monte Carlo with 200 datasets: blue=bias drops, orange=variance rises, "
    "sum = U-shape, minimum at degree 3-4. 'Theory predicts, experiments confirm.' (40s)"
)


# ============================================================
# SLIDE 8: Regularization Motivation
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Regularization: Constraining Complexity",
            font_size=30, bold=True, color=ACCENT)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
            "High variance \u2192 overfitting.  Fix: constrain model complexity.",
            font_size=20, color=TEXT_DARK)

# OLS box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(2.3), Inches(5.5), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

add_textbox(slide, Inches(1.2), Inches(2.5), Inches(4.5), Inches(0.4),
            "Ordinary Least Squares (no penalty):", font_size=16, bold=True, color=TEXT_DARK)
add_latex_image(slide,
    r"\min_{\mathbf{w}} \|\mathbf{y} - \mathbf{Xw}\|_2^2",
    Inches(1.5), Inches(3.2), fontsize=22)

# Ridge box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(7), Inches(2.3), Inches(5.5), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_LIGHT
shape.line.color.rgb = ACCENT

add_textbox(slide, Inches(7.4), Inches(2.5), Inches(4.5), Inches(0.4),
            "Ridge Regression (L2 penalty):", font_size=16, bold=True, color=ACCENT)
add_latex_image(slide,
    r"\min_{\mathbf{w}} \|\mathbf{y} - \mathbf{Xw}\|_2^2 + \alpha\|\mathbf{w}\|_2^2",
    Inches(7.5), Inches(3.2), fontsize=22)

add_block(slide, Inches(2.5), Inches(5.5), Inches(8), Inches(0.85),
          "Course Connection",
          "The L2 norm is the Euclidean norm from Week 3.")

add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "8 / 14", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "[Speaker 2] 'We know WHY overfitting happens \u2014 high variance. How do we FIX it? "
    "By constraining complexity. Ridge adds an L2 penalty \u2014 the Euclidean norm from "
    "Week 3.' (30s)"
)



# ============================================================
# SLIDE 9: Ridge Closed-Form & Interpretations (with visuals)
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Ridge: Closed-Form & Two Interpretations",
            font_size=28, bold=True, color=ACCENT)

# Closed form (compact)
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.3),
            "Closed-form:", font_size=15, bold=True, color=TEXT_DARK)
add_latex_image(slide,
    r"\mathbf{w}_{\mathrm{ridge}} = (\mathbf{X}^\top\mathbf{X} + \alpha\mathbf{I})^{-1}\mathbf{X}^\top\mathbf{y}",
    Inches(2.3), Inches(1.05), fontsize=14)
add_textbox(slide, Inches(9), Inches(1.15), Inches(4), Inches(0.3),
            "vs. OLS: (X'X)\u207b\u00b9 X'y \u2014 only +\u03b1I",
            font_size=12, color=TEXT_GRAY)

# Left panel: Eigenvalue perspective (with diagram)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(1.7), Inches(6.2), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
shape.line.color.rgb = GREEN_ACCENT

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.3),
            "\u2460 Eigenvalue Perspective (Week 5)", font_size=14, bold=True, color=GREEN_ACCENT)

eigen_buf = make_eigenvalue_diagram()
slide.shapes.add_picture(eigen_buf, Inches(0.8), Inches(2.2), width=Inches(5.6))

add_textbox(slide, Inches(0.8), Inches(4.7), Inches(5.8), Inches(0.4),
            "\u03b1I shifts eigenvalues right \u2192 stabilizes inversion \u2192 reduces variance",
            font_size=11, color=TEXT_DARK)

# Right panel: Lagrangian perspective (with contour plot)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(7.0), Inches(1.7), Inches(5.8), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFD, 0xF6, 0xE3)
shape.line.color.rgb = ORANGE

add_textbox(slide, Inches(7.3), Inches(1.8), Inches(5.5), Inches(0.3),
            "\u2461 Lagrangian Perspective (Week 9)", font_size=14, bold=True, color=ORANGE)

contour_buf = make_ridge_constraint_contour()
slide.shapes.add_picture(contour_buf, Inches(8.6), Inches(2.1), width=Inches(3.0))

add_textbox(slide, Inches(7.3), Inches(2.3), Inches(1.5), Inches(2.2),
            "min ||y\u2212Xw||\u00b2\ns.t. ||w||\u00b2 \u2264 t\n\n\u03b1 = Lagrange\nmultiplier",
            font_size=10, color=TEXT_DARK)

add_textbox(slide, Inches(7.3), Inches(4.7), Inches(5.5), Inches(0.4),
            "Larger \u03b1 \u2192 smaller ball \u2192 simpler model",
            font_size=11, color=TEXT_DARK)

# Key insight
add_block(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.85),
          "Key Insight",
          "Adding \u03b1I trades a SMALL bias increase for a LARGE variance decrease.")

add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "9 / 14", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "Ridge closed-form: only difference from OLS is +alphaI.\n"
    "[POINT TO LEFT FIGURE] Eigenvalue view: small eigenvalues make OLS unstable. "
    "Adding alpha shifts all eigenvalues right, stabilizing inversion.\n"
    "[POINT TO RIGHT FIGURE] Lagrangian view: Ridge is constrained optimization. "
    "The circle is ||w||^2 <= t. The ellipses are loss contours. "
    "The tangent point is the Ridge solution.\n"
    "'Small bias increase buys large variance decrease.' (80s)"
)

# ============================================================
# SLIDE 10: Lasso & Geometric Comparison (with L1 vs L2 figure)
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Lasso & L1 vs L2 Geometry",
            font_size=30, bold=True, color=ACCENT)

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(4), Inches(0.4),
            "Lasso (L1 penalty):", font_size=18, bold=True, color=TEXT_DARK)
add_latex_image(slide,
    r"\min_{\mathbf{w}} \|\mathbf{y} - \mathbf{Xw}\|_2^2 + \alpha\|\mathbf{w}\|_1",
    Inches(1.2), Inches(1.7), fontsize=20)

# L1 vs L2 geometry comparison figure (main visual)
geometry_buf = make_l1_vs_l2_geometry()
slide.shapes.add_picture(geometry_buf, Inches(0.3), Inches(2.3), width=Inches(8.0))

# Coefficient paths figure on right
fig_path = os.path.join(FIGURES, "ridge_vs_lasso_coefficients.png")
add_image(slide, fig_path, Inches(8.8), Inches(1.5), width=Inches(4))

add_textbox(slide, Inches(0.3), Inches(5.5), Inches(8), Inches(0.4),
            "L2 circle \u2192 smooth shrinkage   |   L1 diamond \u2192 corners on axes \u2192 sparse solutions",
            font_size=12, bold=True, color=TEXT_DARK)

add_block(slide, Inches(0.8), Inches(5.9), Inches(12), Inches(0.85),
          "Geometric Insight",
          "The L1 ball has corners on coordinate axes. The optimum lands on corners \u2192 "
          "coefficients become exactly zero. Lasso = regularization + feature selection.")

add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "10 / 14", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "[POINT TO LEFT FIGURE] L2 ball = circle, tangent point not on axes "
    "(both w1 and w2 non-zero). L1 ball = diamond with corners ON axes. "
    "The contour hits a corner -> w2=0. That's sparsity!\n"
    "[POINT TO RIGHT FIGURE] Ridge coefficients shrink smoothly. "
    "Lasso coefficients drop to exactly zero one by one.\n"
    "'Lasso = regularization + feature selection.' (60s)"
)
# ============================================================
# SLIDE 11: Cross-Validation
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Choosing \u03b1: K-Fold Cross-Validation",
            font_size=30, bold=True, color=ACCENT)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
            "How to choose the regularization parameter \u03b1?",
            font_size=20, color=TEXT_DARK)

# 4 steps as numbered boxes
steps = [
    ("1", "Partition", "Partition data into K equal folds"),
    ("2", "Train", "Train on K\u22121 folds, evaluate on fold k"),
    ("3", "Average", "Average the K evaluation scores"),
    ("4", "Select", "Select \u03b1 that minimizes CV error"),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.8 + i * 3.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        x, Inches(2.5), Inches(2.8), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_LIGHT if i % 2 == 0 else RGBColor(0xFD, 0xF6, 0xE3)
    shape.line.color.rgb = ACCENT if i % 2 == 0 else ORANGE

    add_textbox(slide, x + Inches(0.2), Inches(2.6), Inches(2.4), Inches(0.5),
                f"Step {num}: {title}", font_size=16, bold=True,
                color=ACCENT if i % 2 == 0 else ORANGE)
    add_textbox(slide, x + Inches(0.2), Inches(3.2), Inches(2.4), Inches(1.2),
                desc, font_size=14, color=TEXT_DARK)

add_block(slide, Inches(2), Inches(5.3), Inches(9), Inches(0.85),
          "Key Advantage",
          "Unbiased test error estimate WITHOUT a separate validation set \u2014 "
          "every data point used for both training and evaluation.")

add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "11 / 14", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "'How to choose alpha? K-fold cross-validation. 4 steps: partition, train, "
    "average, select. Gives unbiased test error estimate without separate "
    "validation set.' (30s)"
)



# ============================================================
# SLIDE 12: Experimental Verification (with figures)
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Experimental Verification: Theory \u2192 Experiment",
            font_size=28, bold=True, color=ACCENT)

# --- Three columns: each prediction + its figure ---
predictions = [
    {
        "title": "\u2460 U-Shaped Test Error",
        "fig": "error_vs_complexity.png",
        "caption": "Test error first drops, then rises\nas polynomial degree increases",
        "color": ACCENT,
        "bg": ACCENT_LIGHT,
    },
    {
        "title": "\u2461 Bias\u2193 Variance\u2191",
        "fig": "bias_variance_tradeoff.png",
        "caption": "Monte Carlo (200 datasets):\nbias drops, variance rises",
        "color": ORANGE,
        "bg": RGBColor(0xFD, 0xF6, 0xE3),
    },
    {
        "title": "\u2462 Ridge Shrinks Coefficients",
        "fig": "ridge_regularization.png",
        "caption": "As \u03b1 increases, coefficients\nshrink and test error drops",
        "color": GREEN_ACCENT,
        "bg": RGBColor(0xE8, 0xF5, 0xE9),
    },
]

col_width = 3.9
gap = 0.3
x_start = 0.5

for idx, pred in enumerate(predictions):
    x = Inches(x_start + idx * (col_width + gap))

    # Colored header bar
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        x, Inches(1.1), Inches(col_width), Inches(0.45))
    header.fill.solid()
    header.fill.fore_color.rgb = pred["color"]
    header.line.fill.background()
    tf = header.text_frame
    tf.paragraphs[0].text = pred["title"]
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Figure
    fig_path = os.path.join(FIGURES, pred["fig"])
    add_image(slide, fig_path, x, Inches(1.7), width=Inches(col_width))

    # Caption below figure
    add_textbox(slide, x + Inches(0.1), Inches(4.8), Inches(col_width - 0.2), Inches(0.7),
                pred["caption"], font_size=11, color=TEXT_DARK,
                alignment=PP_ALIGN.CENTER)

    # Confirmed badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        x + Inches(0.8), Inches(5.6), Inches(2.3), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = pred["bg"]
    badge.line.color.rgb = pred["color"]
    tf = badge.text_frame
    tf.paragraphs[0].text = "\u2705 Confirmed"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = pred["color"]
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Bottom summary
add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
            "All three theoretical predictions confirmed experimentally. Theory predicts, experiments verify.",
            font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "12 / 14", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "[Speaker 3] 'Three predictions, three experiments, all confirmed.\n"
    "[POINT TO LEFT] Prediction 1: test error is U-shaped with complexity. "
    "Our polynomial experiment confirms it.\n"
    "[POINT TO MIDDLE] Prediction 2: bias drops, variance rises. "
    "Monte Carlo with 200 datasets shows exactly this crossover.\n"
    "[POINT TO RIGHT] Prediction 3: Ridge shrinks coefficients. "
    "As alpha increases, weights shrink and test error finds a sweet spot.\n"
    "'Theory predicts, experiments confirm.' (50s)"
)
# ============================================================
# SLIDE 13: Deep Learning + Summary
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Beyond Polynomials & Key Takeaways",
            font_size=30, bold=True, color=ACCENT)

# DL figure on left
fig_path = os.path.join(FIGURES, "dl_synthetic_overfitting.png")
add_image(slide, fig_path, Inches(0.8), Inches(1.3), width=Inches(5.5))
add_textbox(slide, Inches(0.8), Inches(5.2), Inches(5.5), Inches(0.5),
            "Same bias-variance principle applies to neural networks.",
            font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Takeaways on right
add_textbox(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
            "3 Key Takeaways", font_size=22, bold=True, color=ACCENT)

takeaways = [
    ("\u2460", "Bias-Variance Decomposition", "explains the underfitting/overfitting tradeoff mathematically"),
    ("\u2461", "Regularization (Ridge L2 / Lasso L1)", "controls complexity via norm penalties"),
    ("\u2462", "Cross-Validation", "selects the optimal balance point"),
]
for i, (num, title, desc) in enumerate(takeaways):
    y = Inches(2.2 + i * 1.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7), y, Inches(5.8), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = [ACCENT_LIGHT, RGBColor(0xFD,0xF6,0xE3),
                                  RGBColor(0xE8,0xF5,0xE9)][i]
    shape.line.fill.background()

    add_textbox(slide, Inches(7.3), y + Inches(0.1), Inches(5.2), Inches(0.4),
                f"{num}  {title}", font_size=15, bold=True, color=TEXT_DARK)
    add_textbox(slide, Inches(7.3), y + Inches(0.55), Inches(5.2), Inches(0.4),
                desc, font_size=13, color=TEXT_GRAY)

add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "13 / 14", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "'Same principle applies to neural networks \u2014 shallow underfits, deep overfits. "
    "3 takeaways: (1) B-V decomposition explains tradeoff; (2) Regularization controls "
    "complexity; (3) Cross-validation selects balance point.' (40s)"
)


# ============================================================
# SLIDE 14: Interactive Closing — with click animation for answer
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
            "Quick Challenge & Q&A",
            font_size=30, bold=True, color=ACCENT)

# Scenario box (always visible)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.5), Inches(1.3), Inches(10), Inches(2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFD, 0xF6, 0xE3)
shape.line.color.rgb = ORANGE
shape.line.width = Pt(2)

add_textbox(slide, Inches(1.8), Inches(1.5), Inches(9), Inches(0.5),
            "Scenario:", font_size=20, bold=True, color=ORANGE)
add_textbox(slide, Inches(1.8), Inches(2.0), Inches(9), Inches(1),
            "A model achieves 99% training accuracy but only 60% test accuracy.\n\n"
            "Q1: Is this overfitting or underfitting?\n"
            "Q2: What would you recommend to fix it?",
            font_size=18, color=TEXT_DARK)

# --- Click 1: Answer box (hidden initially) ---
ans_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.5), Inches(3.8), Inches(10), Inches(1.5))
ans_box.fill.solid()
ans_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
ans_box.line.color.rgb = GREEN_ACCENT
ans_box.line.width = Pt(2)

ans_label = add_textbox(slide, Inches(1.8), Inches(3.9), Inches(9), Inches(0.4),
            "Answer:", font_size=20, bold=True, color=GREEN_ACCENT)
ans_body = add_textbox(slide, Inches(1.8), Inches(4.4), Inches(9), Inches(0.7),
            "Generalization gap = 39%  \u2192  Overfitting (high variance)\n"
            "\u2192  Add regularization / reduce model complexity / collect more data",
            font_size=17, color=TEXT_DARK)

# Thank you
thankyou = add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.8),
            "Thank You \u2014 Questions?",
            font_size=36, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Add click-to-appear animation for answer
add_appear_animations(slide, [
    [ans_box, ans_label, ans_body, thankyou],  # Click 1: reveal answer + thank you
])

add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
            "14 / 14", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "[INTERACTION] Read scenario. Pause 5s for audience thinking.\n"
    "'Generalization gap is 39% \u2014 clearly overfitting. Fix: add regularization "
    "or reduce complexity.' Thank you. Open for Q&A. (30s)"
)


# ============================================================
# SAVE
# ============================================================
prs.save(OUT)
print(f"Saved to {OUT}")
print(f"Total slides: {len(prs.slides)}")
